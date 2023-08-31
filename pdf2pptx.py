import io
import os
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.util import Cm
from tqdm import trange


def convert_pdf2pptx(
    pdf_file, output_file, resolution, start_page, page_count, quiet=False
):
    doc = fitz.open(pdf_file)
    if not quiet:
        print(pdf_file, "contains", doc.page_count, "slides")

    if page_count is None:
        page_count = doc.page_count

    # transformation matrix: slide to pixmap
    zoom = resolution / 72
    matrix = fitz.Matrix(zoom, zoom, 0)

    # create pptx presentation
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    # configure presentation aspect ratio
    page = doc.load_page(0)
    aspect_ratio = page.rect.width / page.rect.height
    prs.slide_width = int(prs.slide_height * aspect_ratio)

    # create page iterator
    if not quiet:
        page_iter = trange(start_page, start_page + page_count)
    else:
        page_iter = range(start_page, start_page + page_count)

    # iterate over slides
    for page_no in page_iter:
        page = doc.load_page(page_no)

        # write slide as a pixmap
        pixmap = page.get_pixmap(matrix=matrix)
        image_data = pixmap.tobytes(output="JPG")
        image_file = io.BytesIO(image_data)

        # add a slide
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = Cm(0)
        slide.shapes.add_picture(image_file, left, top, height=prs.slide_height)

    if output_file is None:
        output_file = Path(pdf_file).with_suffix(".pptx")

    # save presentation
    prs.save(output_file)


def main():
    """
    Convert a PDF slideshow to Powerpoint PPTX.

    Renders each page as a PNG image and creates the resulting Powerpoint
    slideshow from these images. Useful when you want to use Powerpoint
    to present a set of PDF slides (e.g. slides from Beamer). You can then
    use the presentation capabilities of Powerpoint (notes, ink on slides,
    etc.) with slides created in LaTeX.
    """

    # Get current working directory
    RESOLUTION: int = 300  # image resolution
    PAGE_COUNT = None  # number of pages to convert, None defaults to all slides
    START_PAGE: int = 0  # start page for conversion

    DIR_PATH = Path(os.path.dirname(os.path.abspath(__file__)))
    filenames = os.listdir()
    # print(filenames)

    for file in filenames:
        ext = os.path.splitext(file)[-1].lower()

        if ext == ".pdf":
            convert_pdf2pptx(
                file,
                os.path.splitext(file)[0] + ".pptx",
                RESOLUTION,
                START_PAGE,
                PAGE_COUNT,
                quiet=False,
            )


# if __name__ == "main":
main()
